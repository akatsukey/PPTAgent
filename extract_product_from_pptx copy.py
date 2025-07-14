from dotenv import load_dotenv
import openai
import json
import os
import requests
from pptx import Presentation

# === Load .env ===
load_dotenv()
client = openai.OpenAI()  # Automatically reads API key from env
STRAPI_TOKEN = os.getenv("STRAPI_TOKEN")
STRAPI_URL   = "http://localhost:1337/api/medical-products"   # change host if needed

# === CONFIGURATION ===
PPTX_PATH = "data/WEB MASTER Ver 9.pptx"
OUTPUT_DIR = "products"
MODEL = "gpt-4"  # or gpt-3.5-turbo
POST_TO_STRAPI = False  # Set to True to actually post to Strapi, False to just preview
DEFAULT_IMAGE_ID = 351  # Hardcoded image ID for placeholder
# ======================

def fetch_strapi_data():
    """Fetch real category and division IDs from Strapi API"""
    print("🔄 Fetching category and division data from Strapi...")
    
    # Fetch categories
    categories_response = requests.get("http://localhost:1337/api/categories?pagination[limit]=100")
    if categories_response.status_code != 200:
        print("❌ Failed to fetch categories:", categories_response.status_code)
        return None, None
    
    # Fetch divisions
    divisions_response = requests.get("http://localhost:1337/api/divisions?pagination[limit]=100")
    if divisions_response.status_code != 200:
        print("❌ Failed to fetch divisions:", divisions_response.status_code)
        return None, None
    
    # Parse categories
    categories_data = categories_response.json()
    category_map = {}
    for category in categories_data.get("data", []):
        category_map[category["attributes"]["name"]] = category["id"]
    
    # Parse divisions
    divisions_data = divisions_response.json()
    division_map = {}
    for division in divisions_data.get("data", []):
        division_map[division["attributes"]["name"]] = division["id"]
    
    print(f"✅ Fetched {len(category_map)} categories and {len(division_map)} divisions")
    return category_map, division_map

# Fetch real data from Strapi
category_map, division_map = fetch_strapi_data()

# Fallback to hardcoded values if API fails
if not category_map:
    print("⚠️ Using fallback category map")
    category_map = {
        "Syringes": 59,
        "Needles": 60,
        "Condoms": 58,
        "Scalp vein sets": 61,
        "Blood lancets": 62,
        "Infusion Sets": 63,
        "Infusion Accessories": 64,
        "Transfusion Sets": 65,
        "Extension Tubes": 66,
        "IV CANNULA": 67,
        "Electrodes": 68,
        "Gloves": 69,
        "Adhesive tapes & dressings": 70,
        "Gauze": 71,
        "Disposable Non-Woven Gourments": 72,
        "Suction Catheter": 73,
        "Urinary Catheters": 74,
        "Enteral Feeding Tubes": 75,
        "Hemodialysis Catheters": 76
    }

if not division_map:
    print("⚠️ Using fallback division map")
    division_map = {
        "Infusion": 47,
        "Injectables": 48,
        "Diabetis care": 49,
        "Diagnostic": 50,
        "Hemodialysis": 51,
        "Nutrition": 52,
        "Urology": 53,
        "Wound management": 54,
        "Ortopedic": 55,
        "Protection": 56,
        "Measurement devices": 57,
        "Anesthesia / Respiratory": 58,
        "Surgery": 59,
        "Safety devices": 60,
    }

def extract_text_from_slide(prs, slide_index):
    """Extract all text from a slide"""
    print(f"📄 Extracting text from slide {slide_index + 1}...")
    slide = prs.slides[slide_index]
    text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text.append(shape.text.strip())
    extracted_text = "\n".join(text)
    print(f"✅ Extracted {len(extracted_text)} characters from slide")
    return extracted_text

def build_step1_prompt(slide_text):
    """Build prompt for step 1: General product information"""
    return f"""
Extract general product information from this slide and return ONLY this JSON format:

{{
  "name": "",                ★ Only the product title, do NOT include slide number or division word
  "referenceString": "",     ★ The numeric code at the very start of the title (e.g. 4.1.1)
  "standard": "",            ★ Text after "Manufacturing Standard:" (may be "???" if missing)
  "description": "",         ★ MARKDOWN: 4 subtitles (COMPOSITION, STERILIZATION, PRESENTATIONS, PERFORMANCE) each followed by bullet lines
  "categoryGuess": "",       ★ Choose ONE from {list(category_map.keys())}
  "divisionGuess": ""        ★ Choose ONE from {list(division_map.keys())}
}}

Return ONLY the JSON.

--- Slide Text ---
{slide_text}
--- End ---
"""

def build_step2_prompt(slide_text):
    """Build prompt for step 2: Packaging and table information"""
    return f"""
Extract packaging information and table data from this slide and return ONLY this JSON format:

{{
  "tableInMd": "",           ★ Convert the right-hand table into a Markdown table (first row is the headers you infer)
  "PackagingInformation": {{
    "packing_per_inner_box": 0,        ★ Number of items per inner box (integer)
    "inner_boxes_per_carton": 0,       ★ Number of inner boxes per carton (integer)
    "loading_capacity_20GP": 0,        ★ Loading capacity for 20GP container (integer)
    "loading_capacity_40HC": 0,        ★ Loading capacity for 40HC container (integer)
    "additional_notes": ""             ★ Any additional packaging notes (string)
  }}
}}

If packaging information is not available, set all packaging values to 0 and additional_notes to "Not specified".
If no table is present, set tableInMd to "".

Return ONLY the JSON.

--- Slide Text ---
{slide_text}
--- End ---
"""

def call_openai(prompt, step_name):
    """Call OpenAI API with logging"""
    print(f"🔁 Step {step_name}: Sending request to OpenAI...")
    response = client.chat.completions.create(
        model=MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3
    )
    print(f"✅ Step {step_name}: Received response from OpenAI.")
    return response.choices[0].message.content

def merge_step_results(step1_result, step2_result):
    """Merge results from both steps into final Strapi format"""
    print("🔄 Merging step results...")
    
    # Extract data from step 1
    step1_data = json.loads(step1_result)
    
    # Extract data from step 2
    step2_data = json.loads(step2_result)
    
    # Build final payload
    final_payload = {
        "data": {
            "name": step1_data.get("name"),
            "referenceString": step1_data.get("referenceString"),
            "standard": step1_data.get("standard", "???"),
            "description": step1_data.get("description", ""),
            "tableInMd": step2_data.get("tableInMd", ""),
            "category": category_map.get(step1_data.get("categoryGuess")),
            "divisions": [division_map.get(step1_data.get("divisionGuess"))] if step1_data.get("divisionGuess") else [],
            "images": [],
            "PackagingInformation": step2_data.get("PackagingInformation")
        }
    }
    
    print("✅ Results merged successfully")
    return final_payload

def save_json(payload, slide_index):
    """Save the final payload to a JSON file"""
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    # Create filename based on slide number
    filename = os.path.join(OUTPUT_DIR, f"slide_{slide_index + 1}_strapi_format.json")
    
    with open(filename, "w", encoding="utf-8") as f:
        json.dump(payload, f, indent=2, ensure_ascii=False)
    
    print(f"📁 Saved to {filename}")
    return filename

def process_slide_two_steps(prs, slide_index):
    """Process a slide using the two-step approach"""
    print(f"\n🎯 Processing slide {slide_index + 1} using two-step extraction...")
    
    # Step 1: Extract general product information
    print("\n📋 STEP 1: Extracting general product information...")
    slide_text = extract_text_from_slide(prs, slide_index)
    step1_prompt = build_step1_prompt(slide_text)
    step1_result = call_openai(step1_prompt, "1")
    print("✅ Step 1 completed")
    
    # Step 2: Extract packaging and table information
    print("\n📋 STEP 2: Extracting packaging and table information...")
    step2_prompt = build_step2_prompt(slide_text)
    step2_result = call_openai(step2_prompt, "2")
    print("✅ Step 2 completed")
    
    # Merge results
    final_payload = merge_step_results(step1_result, step2_result)
    
    # Save to file
    filename = save_json(final_payload, slide_index)
    
    # Display results
    print(f"\n📊 Final payload for slide {slide_index + 1}:")
    print(json.dumps(final_payload, indent=2))
    
    return final_payload

def main():
    print("🚀 Starting two-step product extraction from PowerPoint...")
    print(f"📤 POST_TO_STRAPI: {POST_TO_STRAPI}")
    print()
    
    prs = Presentation(PPTX_PATH)
    print(f"📄 Loaded presentation with {len(prs.slides)} slides")

    while True:
        try:
            slide_input = input(f"\n📌 Enter slide number to process (1–{len(prs.slides)}) or 'q' to quit: ")
            
            if slide_input.lower() == "q":
                print("👋 Exiting.")
                break
            
            slide_num = int(slide_input) - 1  # Adjust to 0-based index
            
            if 0 <= slide_num < len(prs.slides):
                process_slide_two_steps(prs, slide_num)
                
                # Ask if user wants to process another slide
                continue_input = input(f"\n🔄 Process another slide? (y/n): ")
                if continue_input.lower() != 'y':
                    print("👋 Exiting.")
                    break
            else:
                print("❌ Invalid slide number. Please try again.")
                
        except ValueError:
            print("❌ Please enter a valid number or 'q'.")
        except KeyboardInterrupt:
            print("\n👋 Exiting.")
            break
        except Exception as e:
            print(f"❌ Error processing slide: {e}")
            print("Please try again.")

if __name__ == "__main__":
    main()
