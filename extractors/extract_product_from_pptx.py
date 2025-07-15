from dotenv import load_dotenv
import openai
import json
import os
import requests
from pptx import Presentation
import sys
import argparse

# === Load .env ===
load_dotenv()
client = openai.OpenAI()  # Automatically reads API key from env
STRAPI_TOKEN = os.getenv("STRAPI_TOKEN")
print("strapi token", STRAPI_TOKEN)
STRAPI_URL   = "http://localhost:1337/api/medical-products"   # change host if needed

# === CONFIGURATION ===
PPTX_PATH = "../data/WEB MASTER Ver 9.pptx"
OUTPUT_DIR = "products"
MODEL = "gpt-4"  # or gpt-3.5-turbo
POST_TO_STRAPI = True  # Set to True to actually post to Strapi, False to just preview
DEFAULT_IMAGE_ID = 351  # Hardcoded image ID for placeholder
HARDCODED_IMAGE_ID = 351  # Image ID to use when images array is empty
# ======================

# === COMMAND LINE ARGUMENT PARSING ===
def parse_arguments():
    """Parse command line arguments for slide range"""
    parser = argparse.ArgumentParser(
        description="Extract product information from PowerPoint slides",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python extract_product_from_pptx.py --start 11 --finish 15
  python extract_product_from_pptx.py -s 21 -f 28
  python extract_product_from_pptx.py --start 1 --finish 10
        """
    )
    
    parser.add_argument(
        "--start", "-s",
        type=int,
        required=True,
        help="Starting slide number (1-based)"
    )
    
    parser.add_argument(
        "--finish", "-f", 
        type=int,
        required=True,
        help="Finishing slide number (1-based)"
    )
    
    parser.add_argument(
        "--no-update",
        action="store_true",
        help="Skip updating products.json file"
    )
    
    return parser.parse_args()

# === TASK PLAN ===
"""
🎯 TASK PLAN: Fixing Extraction Issues

A) TABLE EXTRACTION (python-pptx)
   - Iterate through slide.shapes
   - Check shape.has_table for each shape
   - Extract table data (rows, columns, cell values)
   - Convert to Markdown table format
   - Handle edge cases (empty tables, merged cells)

B) MARKDOWN TABLE CONVERSION
   - Parse table structure (headers, data rows)
   - Generate proper Markdown table syntax
   - Handle special characters and formatting
   - Ensure proper alignment and spacing

C) REFERENCESTRING PADDING
   - Parse referenceString (e.g., "4.1.1")
   - Split by dots into segments
   - Zero-pad middle segment to 2 digits
   - Convert "4.1.1" → "4.01.1"

D) DESCRIPTION FORMATTING
   - Parse raw description sections
   - Add ## headings for each section
   - Convert bullet points to proper Markdown
   - Ensure proper spacing and structure

E) IMAGES ARRAY HANDLING
   - Check if images array is empty
   - Insert HARDCODED_IMAGE_ID if needed
   - Maintain array structure for Strapi

F) COMPLETENESS LOGGING
   - Track field completion status
   - Generate ✅/❌ indicators
   - Log missing or incomplete fields
   - Provide actionable feedback
"""

# === STUB FUNCTIONS (NOT YET INTEGRATED) ===

def extract_table_md(slide):
    """
    Extract table from slide and convert to Markdown format.
    
    Args:
        slide: pptx.slide.Slide object
        
    Returns:
        str: Markdown table string or empty string if no table found
    """
    print("🔍 Searching for tables in slide...")
    
    for shape in slide.shapes:
        if shape.has_table:
            print("📊 Found table, converting to Markdown...")
            table = shape.table
            
            # Extract table data
            rows = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    # Get cell text and clean it
                    cell_text = cell.text.strip()
                    # Replace pipe characters to avoid breaking markdown
                    cell_text = cell_text.replace('|', '\\|')
                    row_data.append(cell_text)
                rows.append(row_data)
            
            if not rows:
                print("⚠️ Table found but no data extracted")
                return ""
            
            # Convert to Markdown format
            markdown_table = ""
            
            # Add header row
            if rows:
                markdown_table += "| " + " | ".join(rows[0]) + " |\n"
                # Add separator row
                markdown_table += "|" + "|".join(["---"] * len(rows[0])) + "|\n"
                # Add data rows
                for row in rows[1:]:
                    markdown_table += "| " + " | ".join(row) + " |\n"
            
            print(f"✅ Converted table with {len(rows)} rows to Markdown")
            return markdown_table.strip()
    
    print("❌ No table found in slide")
    return ""

def format_description(raw_sections_dict):
    """
    Format raw description sections into proper Markdown.
    
    Args:
        raw_sections_dict: dict with sections like {"COMPOSITION": [...], "STERILIZATION": [...]}
        
    Returns:
        str: Properly formatted Markdown string with ## headings and bullets
    """
    print("📝 Formatting description with Markdown structure...")
    
    # TODO: Parse raw sections
    # TODO: Add ## headings for each section
    # TODO: Convert bullet points to proper Markdown
    # TODO: Handle missing sections gracefully
    
    formatted_md = ""
    for section_name, bullet_points in raw_sections_dict.items():
        formatted_md += f"## {section_name}\n"
        for point in bullet_points:
            formatted_md += f"- {point}\n"
        formatted_md += "\n"
    
    return formatted_md.strip()

def pad_reference_string(reference_string):
    """
    Pad referenceString to ensure middle segment has 2 digits.
    
    Args:
        reference_string: str like "4.1.1" or "4.01.1"
        
    Returns:
        str: Padded reference string like "4.01.1"
    """
    print(f"🔢 Padding reference string: {reference_string}")
    
    # TODO: Parse reference string by dots
    # TODO: Zero-pad middle segment if needed
    # TODO: Handle edge cases (invalid formats)
    
    if not reference_string:
        return ""
    
    parts = reference_string.split('.')
    if len(parts) >= 3:
        # Ensure middle segment has 2 digits
        parts[1] = parts[1].zfill(2)
        return '.'.join(parts)
    
    return reference_string

def ensure_images_array(payload_data):
    """
    Ensure images array is not empty by adding HARDCODED_IMAGE_ID.
    
    Args:
        payload_data: dict containing the data payload
        
    Returns:
        dict: Updated payload with images array populated if needed
    """
    print("🖼️ Checking images array...")
    
    if "images" not in payload_data or not payload_data["images"]:
        print(f"✅ Adding HARDCODED_IMAGE_ID ({HARDCODED_IMAGE_ID}) to images array")
        payload_data["images"] = [HARDCODED_IMAGE_ID]
    else:
        print(f"✅ Images array already populated with {len(payload_data['images'])} items")
    
    return payload_data

def generate_completeness_log(payload_data):
    """
    Generate completeness log with ✅/❌ indicators for each field.
    
    Args:
        payload_data: dict containing the data payload
        
    Returns:
        dict: Field status with completion indicators
    """
    print("📋 Generating completeness log...")
    
    completeness_log = {}
    
    # Define field requirements
    required_fields = {
        "name": "Product name",
        "referenceString": "Product reference code",
        "images": "Images array"
    }
    
    optional_fields = {
        "standard": "Manufacturing standard",
        "description": "Product description (markdown)",
        "tableInMd": "Technical specifications table",
        "category": "Product category ID",
        "divisions": "Product division IDs",
        "PackagingInformation": "Packaging component"
    }
    
    # Check required fields
    for field, description in required_fields.items():
        if field in payload_data and payload_data[field]:
            if field == "images" and len(payload_data[field]) > 0:
                completeness_log[field] = ("✅", f"{description} - {len(payload_data[field])} image(s)")
            else:
                completeness_log[field] = ("✅", f"{description} - populated")
        else:
            completeness_log[field] = ("❌", f"{description} - missing or empty")
    
    # Check optional fields
    for field, description in optional_fields.items():
        if field in payload_data and payload_data[field]:
            if field == "description":
                preview = payload_data[field][:50] + "..." if len(payload_data[field]) > 50 else payload_data[field]
                completeness_log[field] = ("✅", f"{description} - \"{preview}\"")
            elif field == "tableInMd":
                if payload_data[field]:
                    completeness_log[field] = ("✅", f"{description} - table present")
                else:
                    completeness_log[field] = ("❌", f"{description} - empty")
            else:
                completeness_log[field] = ("✅", f"{description} - populated")
        else:
            completeness_log[field] = ("❌", f"{description} - missing")
    
    return completeness_log

def print_completeness_log(completeness_log):
    """
    Print the completeness log with visual indicators.
    
    Args:
        completeness_log: dict from generate_completeness_log()
    """
    print("\n📊 COMPLETENESS REPORT:")
    print("=" * 50)
    
    for field, (status, description) in completeness_log.items():
        print(f"{status} {field}: {description}")
    
    # Summary
    total_fields = len(completeness_log)
    completed_fields = sum(1 for status, _ in completeness_log.values() if status == "✅")
    completion_rate = (completed_fields / total_fields) * 100
    
    print("=" * 50)
    print(f"📈 Completion Rate: {completed_fields}/{total_fields} ({completion_rate:.1f}%)")
    
    if completion_rate < 80:
        print("⚠️  Warning: Low completion rate - consider manual review")
    elif completion_rate == 100:
        print("🎉 Perfect! All fields completed successfully")
    else:
        print("✅ Good completion rate")

# === END STUB FUNCTIONS ===

def fetch_strapi_data():
    """Fetch real category and division IDs from Strapi API"""
    print("🔄 Fetching category and division data from Strapi...")
    
    # Fetch categories
    categories_response = requests.get("http://localhost:1337/api/categories?pagination[limit]=100")
    if categories_response.status_code != 200:
        print("❌ Failed to fetch categories:", categories_response.status_code)
        return None, None
    
    # Fetch divisions
    divisions_response = requests.get("http://localhost:1337/api/divisions?pagination[limit]=100")
    if divisions_response.status_code != 200:
        print("❌ Failed to fetch divisions:", divisions_response.status_code)
        return None, None
    
    # Parse categories
    categories_data = categories_response.json()
    category_map = {}
    for category in categories_data.get("data", []):
        category_map[category["attributes"]["name"]] = category["id"]
    
    # Parse divisions
    divisions_data = divisions_response.json()
    division_map = {}
    for division in divisions_data.get("data", []):
        division_map[division["attributes"]["name"]] = division["id"]
    
    print(f"✅ Fetched {len(category_map)} categories and {len(division_map)} divisions")
    return category_map, division_map

# Fetch real data from Strapi
category_map, division_map = fetch_strapi_data()

# Fallback to hardcoded values if API fails
if not category_map:
    print("⚠️ Using fallback category map")
    category_map = {
        "Syringes": 59,
        "Needles": 60,
        "Condoms": 58,
        "Scalp vein sets": 61,
        "Blood lancets": 62,
        "Infusion Sets": 63,
        "Infusion Accessories": 64,
        "Transfusion Sets": 65,
        "Extension Tubes": 66,
        "IV CANNULA": 67,
        "Electrodes": 68,
        "Gloves": 69,
        "Adhesive tapes & dressings": 70,
        "Gauze": 71,
        "Disposable Non-Woven Gourments": 72,
        "Suction Catheter": 73,
        "Urinary Catheters": 74,
        "Enteral Feeding Tubes": 75,
        "Hemodialysis Catheters": 76
    }

if not division_map:
    print("⚠️ Using fallback division map")
    division_map = {
        "Infusion": 47,
        "Injectables": 48,
        "Diabetis care": 49,
        "Diagnostic": 50,
        "Hemodialysis": 51,
        "Nutrition": 52,
        "Urology": 53,
        "Wound management": 54,
        "Ortopedic": 55,
        "Protection": 56,
        "Measurement devices": 57,
        "Anesthesia / Respiratory": 58,
        "Surgery": 59,
        "Safety devices": 60,
    }

def extract_text_from_slide(prs, slide_index):
    """Extract all text from a slide"""
    print(f"📄 Extracting text from slide {slide_index + 1}...")
    slide = prs.slides[slide_index]
    text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text.append(shape.text.strip())
    extracted_text = "\n".join(text)
    print(f"✅ Extracted {len(extracted_text)} characters from slide")
    return extracted_text

def build_step1_prompt(slide_text):
    """Build prompt for step 1: General product information"""
    return f"""
Extract general product information from this slide and return ONLY this JSON format:

{{
  "name": "",                ★ Only the product title, do NOT include slide number or division word
  "referenceString": "",     ★ The numeric code at the very start of the title (e.g. 4.1.1)
  "standard": "",            ★ Text after "Manufacturing Standard:" (may be "???" if missing)
  "description": "",         ★ MARKDOWN FORMAT: Use ## for section headers and - for bullet points. Example:
                             ★ ## COMPOSITION
                             ★ - Barrel + Plunger: Medical Grade Polypropylene (PP)
                             ★ - Gasket: Isoprene Latex-Free Rubber and Natural Rubber
                             ★ 
                             ★ ## STERILIZATION
                             ★ - Ethylene Oxide (EO) Sterile
  "categoryGuess": "",       ★ Choose ONE from {list(category_map.keys())}
  "divisionGuess": ""        ★ Choose ONE from {list(division_map.keys())}
}}

IMPORTANT: For the description field, format each section with ## headers and - bullet points exactly as shown in the example above.

Return ONLY the JSON.

--- Slide Text ---
{slide_text}
--- End ---
"""

def build_step2_prompt(slide_text):
    """Build prompt for step 2: Packaging information"""
    return f"""
Extract packaging information from this slide and return ONLY this JSON format:

{{
  "tableInMd": "",           ★ Leave empty - table extraction is handled separately
  "PackagingInformation": {{
    "packing_per_inner_box": 0,        ★ Number of items per inner box (integer)
    "inner_boxes_per_carton": 0,       ★ Number of inner boxes per carton (integer)
    "loading_capacity_20GP": 0,        ★ Loading capacity for 20GP container (integer)
    "loading_capacity_40HC": 0,        ★ Loading capacity for 40HC container (integer)
    "additional_notes": ""             ★ Any additional packaging notes (string)
  }}
}}

If packaging information is not available, set all packaging values to 0 and additional_notes to "Not specified".

Return ONLY the JSON.

--- Slide Text ---
{slide_text}
--- End ---
"""

def call_openai(prompt, step_name):
    """Call OpenAI API with logging"""
    print(f"🔁 Step {step_name}: Sending request to OpenAI...")
    response = client.chat.completions.create(
        model=MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3
    )
    print(f"✅ Step {step_name}: Received response from OpenAI.")
    return response.choices[0].message.content

def merge_step_results(step1_result, step2_result):
    """Merge results from both steps into final Strapi format"""
    print("🔄 Merging step results...")
    
    # Extract data from step 1
    step1_data = json.loads(step1_result)
    
    # Extract data from step 2
    step2_data = json.loads(step2_result)
    
    # Build final payload
    final_payload = {
        "data": {
            "name": step1_data.get("name"),
            "referenceString": step1_data.get("referenceString"),
            "standard": step1_data.get("standard", "???"),
            "description": step1_data.get("description", ""),
            "tableInMd": step2_data.get("tableInMd", ""),
            "category": category_map.get(step1_data.get("categoryGuess")),
            "divisions": [division_map.get(step1_data.get("divisionGuess"))] if step1_data.get("divisionGuess") else [],
            "images": [],
            "PackagingInformation": step2_data.get("PackagingInformation")
        }
    }
    
    print("✅ Results merged successfully")
    return final_payload

def save_json(payload, slide_index):
    """Save the final payload to a JSON file"""
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    # Create filename based on slide number
    filename = os.path.join(OUTPUT_DIR, f"slide_{slide_index + 1}_strapi_format.json")
    
    with open(filename, "w", encoding="utf-8") as f:
        json.dump(payload, f, indent=2, ensure_ascii=False)
    
    print(f"📁 Saved to {filename}")
    return filename

def process_slide_two_steps(prs, slide_index):
    """Process a slide using the two-step approach"""
    print(f"\n🎯 Processing slide {slide_index + 1} using two-step extraction...")
    
    # Get the slide object for table extraction
    slide = prs.slides[slide_index]
    
    # Step 1: Extract general product information using GPT-4
    print("\n📋 STEP 1: Extracting general product information...")
    slide_text = extract_text_from_slide(prs, slide_index)
    step1_prompt = build_step1_prompt(slide_text)
    step1_result = call_openai(step1_prompt, "1")
    print("✅ Step 1 completed")
    
    # Step 2: Extract packaging information using GPT-4 and table using python-pptx
    print("\n📊 STEP 2: Extracting packaging information and table data...")
    
    # Extract table using python-pptx
    table_md = extract_table_md(slide)
    
    # Extract packaging information using GPT-4
    step2_prompt = build_step2_prompt(slide_text)
    step2_result = call_openai(step2_prompt, "2")
    print("✅ Step 2 completed")
    
    # Parse step 2 result and replace tableInMd with actual extracted table
    step2_data = json.loads(step2_result)
    if table_md:
        step2_data["tableInMd"] = table_md
        print("✅ Replaced GPT-extracted table with actual table data")
    
    # Merge results
    final_payload = merge_step_results(step1_result, json.dumps(step2_data))
    
    # Save to file
    filename = save_json(final_payload, slide_index)
    
    # Display results
    print(f"\n📊 Final payload for slide {slide_index + 1}:")
    print(json.dumps(final_payload, indent=2))
    
    return final_payload

def update_products_json(slide_payloads, products_json_path="../products.json"):
    """Update products.json with extracted slide data."""
    print(f"\n🔄 Updating {products_json_path} with extracted slide data...")
    try:
        with open(products_json_path, "r", encoding="utf-8") as f:
            products = json.load(f)
    except Exception as e:
        print(f"❌ Could not load {products_json_path}: {e}")
        return

    ref_to_product = {p["data"]["referenceString"]: p for p in products}
    updated_count = 0
    for payload in slide_payloads:
        ref = payload["data"].get("referenceString")
        if not ref:
            print(f"⚠️  No referenceString in slide payload, skipping.")
            continue
        if ref not in ref_to_product:
            print(f"⚠️  No matching product for referenceString {ref}, skipping.")
            continue
        prod = ref_to_product[ref]["data"]
        for key in ["description", "standard", "tableInMd", "PackagingInformation"]:
            new_val = payload["data"].get(key)
            if new_val and new_val != "???" and new_val != "??????????":
                prod[key] = new_val
        updated_count += 1
        print(f"✅ Updated product {prod['name']} (Ref: {ref})")
    # Write back to file
    with open(products_json_path, "w", encoding="utf-8") as f:
        json.dump(products, f, indent=2, ensure_ascii=False)
    print(f"\n✅ Updated {updated_count} products in {products_json_path}")

def main():
    # Parse command line arguments
    args = parse_arguments()
    
    # Convert 1-based slide numbers to 0-based indices
    start_slide = args.start - 1  # Convert to 0-based index
    end_slide = args.finish - 1   # Convert to 0-based index
    
    print("🚀 Starting two-step product extraction from PowerPoint...")
    print(f"📤 POST_TO_STRAPI: {POST_TO_STRAPI}")
    print(f"🎯 Processing slides {args.start} to {args.finish} (indices {start_slide} to {end_slide})")
    print()
    
    # Validate slide range
    if start_slide < 0:
        print(f"❌ Error: Start slide {args.start} is invalid (must be >= 1)")
        sys.exit(1)
    
    if end_slide < start_slide:
        print(f"❌ Error: End slide {args.finish} must be >= start slide {args.start}")
        sys.exit(1)
    
    prs = Presentation(PPTX_PATH)
    print(f"📄 Loaded presentation with {len(prs.slides)} slides")
    
    # Check if slide range is valid
    if end_slide >= len(prs.slides):
        print(f"❌ Error: End slide {args.finish} exceeds presentation length ({len(prs.slides)} slides)")
        sys.exit(1)
    
    all_results = []
    
    for slide_index in range(start_slide, end_slide + 1):
        try:
            print(f"\n{'='*60}")
            print(f"📄 PROCESSING SLIDE {slide_index + 1}")
            print(f"{'='*60}")
            
            result = process_slide_two_steps(prs, slide_index)
            all_results.append(result)
            
            print(f"✅ Completed slide {slide_index + 1}")
            
        except Exception as e:
            print(f"❌ Error processing slide {slide_index + 1}: {e}")
            continue
    
    # Save all results to a combined file
    if all_results:
        combined_filename = os.path.join(OUTPUT_DIR, f"slides_{args.start}_to_{args.finish}_combined.json")
        with open(combined_filename, "w", encoding="utf-8") as f:
            json.dump(all_results, f, indent=2, ensure_ascii=False)
        print(f"\n📁 Saved combined results to {combined_filename}")
        print(f"✅ Successfully processed {len(all_results)} slides")
        
        # Update products.json unless --no-update flag is used
        if not args.no_update:
            update_products_json(all_results, products_json_path="../products.json")
        else:
            print("⏭️  Skipping products.json update (--no-update flag used)")
    else:
        print("\n❌ No slides were successfully processed")

if __name__ == "__main__":
    main()
