#!/usr/bin/env python3
"""
Slide Content Extractor - Extracts all text content from PowerPoint slides
in their original format, including tables, shapes, and text boxes.
"""

import argparse
import json
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# === CONFIGURATION ===
PPTX_PATH = "../data/WEB MASTER Ver 9.pptx"
OUTPUT_DIR = "slide_content"

def parse_arguments():
    """Parse command line arguments for slide range"""
    parser = argparse.ArgumentParser(
        description="Extract all text content from PowerPoint slides in original format",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python extract_slide_content.py --start 11 --finish 15
  python extract_slide_content.py -s 21 -f 28
  python extract_slide_content.py --start 1 --finish 10
        """
    )
    
    parser.add_argument(
        "--start", "-s",
        type=int,
        required=True,
        help="Starting slide number (1-based)"
    )
    
    parser.add_argument(
        "--finish", "-f", 
        type=int,
        required=True,
        help="Finishing slide number (1-based)"
    )
    
    parser.add_argument(
        "--output", "-o",
        type=str,
        default="slide_content.json",
        help="Output JSON file name"
    )
    
    return parser.parse_args()

def extract_text_from_shape(shape):
    """Extract text from a shape, handling different shape types"""
    text_content = {
        "type": "unknown",
        "text": "",
        "position": None,
        "size": None
    }
    
    # Get shape type
    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        text_content["type"] = "text_box"
        if shape.has_text_frame:
            text_content["text"] = shape.text.strip()
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        text_content["type"] = "table"
        text_content["text"] = extract_table_text(shape.table)
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        text_content["type"] = "picture"
        text_content["text"] = f"[Image: {shape.name if shape.name else 'unnamed'}]"
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        text_content["type"] = "group"
        text_content["text"] = "[Group of shapes]"
    else:
        text_content["type"] = f"shape_{shape.shape_type}"
        if hasattr(shape, 'text') and shape.text:
            text_content["text"] = shape.text.strip()
    
    # Get position and size
    if hasattr(shape, 'left') and hasattr(shape, 'top'):
        text_content["position"] = {
            "left": shape.left,
            "top": shape.top
        }
    
    if hasattr(shape, 'width') and hasattr(shape, 'height'):
        text_content["size"] = {
            "width": shape.width,
            "height": shape.height
        }
    
    return text_content

def extract_table_text(table):
    """Extract text from a table in a structured format"""
    table_data = {
        "rows": [],
        "headers": [],
        "markdown": ""
    }
    
    if not table.rows:
        return table_data
    
    # Extract headers (first row)
    headers = []
    for cell in table.rows[0].cells:
        headers.append(cell.text.strip())
    table_data["headers"] = headers
    
    # Extract all rows
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        table_data["rows"].append(row_data)
    
    # Generate Markdown table
    if headers and table_data["rows"]:
        markdown_table = "| " + " | ".join(headers) + " |\n"
        markdown_table += "| " + " | ".join(["---"] * len(headers)) + " |\n"
        
        for row in table_data["rows"][1:]:  # Skip header row
            markdown_table += "| " + " | ".join(row) + " |\n"
        
        table_data["markdown"] = markdown_table.strip()
    
    return table_data

def extract_slide_content(prs, slide_index):
    """Extract all content from a slide in structured format"""
    slide = prs.slides[slide_index]
    slide_content = {
        "slide_number": slide_index + 1,
        "shapes": [],
        "raw_text": "",
        "tables": [],
        "text_boxes": [],
        "images": [],
        "other_shapes": []
    }
    
    all_text = []
    
    for shape in slide.shapes:
        shape_content = extract_text_from_shape(shape)
        slide_content["shapes"].append(shape_content)
        
        # Categorize shapes
        if shape_content["type"] == "table":
            slide_content["tables"].append(shape_content)
        elif shape_content["type"] == "text_box":
            slide_content["text_boxes"].append(shape_content)
        elif shape_content["type"] == "picture":
            slide_content["images"].append(shape_content)
        else:
            slide_content["other_shapes"].append(shape_content)
        
        # Collect all text
        if shape_content["text"]:
            if isinstance(shape_content["text"], dict):
                # Handle table content
                if "markdown" in shape_content["text"] and shape_content["text"]["markdown"]:
                    all_text.append(shape_content["text"]["markdown"])
                elif "rows" in shape_content["text"]:
                    # Convert table rows to text
                    table_text = []
                    for row in shape_content["text"]["rows"]:
                        table_text.append(" | ".join(row))
                    all_text.append("\n".join(table_text))
            else:
                all_text.append(shape_content["text"])
    
    # Combine all text
    slide_content["raw_text"] = "\n\n".join(all_text)
    
    return slide_content

def save_slide_content(slide_content, slide_index, output_dir):
    """Save slide content to individual JSON file"""
    os.makedirs(output_dir, exist_ok=True)
    
    filename = os.path.join(output_dir, f"slide_{slide_index + 1}_content.json")
    
    with open(filename, "w", encoding="utf-8") as f:
        json.dump(slide_content, f, indent=2, ensure_ascii=False)
    
    print(f"📁 Saved slide {slide_index + 1} content to {filename}")
    return filename

def print_slide_summary(slide_content):
    """Print a summary of what was extracted from the slide"""
    print(f"\n📊 SLIDE {slide_content['slide_number']} SUMMARY:")
    print(f"   📝 Text boxes: {len(slide_content['text_boxes'])}")
    print(f"   📊 Tables: {len(slide_content['tables'])}")
    print(f"   🖼️  Images: {len(slide_content['images'])}")
    print(f"   🔲 Other shapes: {len(slide_content['other_shapes'])}")
    
    # Show table previews
    for i, table in enumerate(slide_content['tables']):
        print(f"   📋 Table {i+1}: {len(table['rows'])} rows")
        if table['markdown']:
            print(f"      Preview: {table['markdown'][:100]}...")
    
    # Show text previews
    for i, text_box in enumerate(slide_content['text_boxes']):
        preview = text_box['text'][:50] + "..." if len(text_box['text']) > 50 else text_box['text']
        print(f"   📝 Text box {i+1}: {preview}")

def main():
    """Main function to extract slide content"""
    args = parse_arguments()
    
    # Convert 1-based slide numbers to 0-based indices
    start_slide = args.start - 1
    end_slide = args.finish - 1
    
    print("🔍 Starting slide content extraction...")
    print(f"🎯 Processing slides {args.start} to {args.finish} (indices {start_slide} to {end_slide})")
    print()
    
    # Validate slide range
    if start_slide < 0:
        print(f"❌ Error: Start slide {args.start} is invalid (must be >= 1)")
        sys.exit(1)
    
    if end_slide < start_slide:
        print(f"❌ Error: End slide {args.finish} must be >= start slide {args.start}")
        sys.exit(1)
    
    # Load presentation
    try:
        prs = Presentation(PPTX_PATH)
        print(f"📄 Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"❌ Error loading presentation: {e}")
        sys.exit(1)
    
    # Check if slide range is valid
    if end_slide >= len(prs.slides):
        print(f"❌ Error: End slide {args.finish} exceeds presentation length ({len(prs.slides)} slides)")
        sys.exit(1)
    
    all_slide_content = []
    
    for slide_index in range(start_slide, end_slide + 1):
        try:
            print(f"\n{'='*60}")
            print(f"📄 EXTRACTING SLIDE {slide_index + 1}")
            print(f"{'='*60}")
            
            slide_content = extract_slide_content(prs, slide_index)
            all_slide_content.append(slide_content)
            
            # Save individual slide content
            save_slide_content(slide_content, slide_index, OUTPUT_DIR)
            
            # Print summary
            print_slide_summary(slide_content)
            
            print(f"✅ Completed slide {slide_index + 1}")
            
        except Exception as e:
            print(f"❌ Error processing slide {slide_index + 1}: {e}")
            continue
    
    # Save combined results
    if all_slide_content:
        combined_filename = os.path.join(OUTPUT_DIR, f"slides_{args.start}_to_{args.finish}_combined.json")
        with open(combined_filename, "w", encoding="utf-8") as f:
            json.dump(all_slide_content, f, indent=2, ensure_ascii=False)
        print(f"\n📁 Saved combined results to {combined_filename}")
        print(f"✅ Successfully processed {len(all_slide_content)} slides")
        
        # Print overall summary
        total_tables = sum(len(slide['tables']) for slide in all_slide_content)
        total_text_boxes = sum(len(slide['text_boxes']) for slide in all_slide_content)
        total_images = sum(len(slide['images']) for slide in all_slide_content)
        
        print(f"\n📊 OVERALL SUMMARY:")
        print(f"   📄 Slides processed: {len(all_slide_content)}")
        print(f"   📊 Total tables found: {total_tables}")
        print(f"   📝 Total text boxes: {total_text_boxes}")
        print(f"   🖼️  Total images: {total_images}")
        
    else:
        print("\n❌ No slides were successfully processed")

if __name__ == "__main__":
    main() 