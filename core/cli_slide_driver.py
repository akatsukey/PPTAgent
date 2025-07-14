#!/usr/bin/env python3
"""
Interactive CLI for processing PowerPoint slides with human-in-the-loop decision making.

This script provides an interactive interface for stepping through slides in a PowerPoint
presentation and deciding whether to extract, diff, upload, delete-and-recreate, or skip
each slide. It integrates with the existing extractor and Strapi API, with resumable
progress tracking.
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass, asdict
from datetime import datetime
import requests
from dotenv import load_dotenv
from pptx import Presentation
from rich.console import Console
from rich.panel import Panel
from rich.table import Table
from rich.prompt import Prompt, Confirm
from rich.text import Text
from slugify import slugify

# Import existing extractor
from extract_product_from_pptx import process_slide_two_steps, fetch_strapi_data

# Load environment variables
load_dotenv()

# Initialize Rich console
console = Console()

@dataclass
class ProgressEntry:
    """Represents progress for a single slide"""
    slide_number: int
    action: str  # 'extract', 'upload', 'delete', 'skip', 'diff'
    timestamp: str
    product_name: Optional[str] = None
    strapi_id: Optional[int] = None
    error: Optional[str] = None

@dataclass
class SlideInfo:
    """Information about a slide"""
    slide_number: int
    product_name: str
    reference_string: str
    existing_in_strapi: bool
    strapi_id: Optional[int] = None
    strapi_url: Optional[str] = None

class StrapiUtils:
    """Utility class for Strapi API operations"""
    
    def __init__(self, base_url: str = "http://localhost:1337", token: Optional[str] = None):
        self.base_url = base_url
        self.token = token or os.getenv("STRAPI_TOKEN")
        self.headers = {
            "Authorization": f"Bearer {self.token}",
            "Content-Type": "application/json"
        } if self.token else {}
    
    def find_by_name(self, name: str) -> Optional[Dict]:
        """Find a product by name in Strapi"""
        try:
            response = requests.get(
                f"{self.base_url}/api/medical-products",
                params={"filters[name][$eq]": name},
                headers=self.headers
            )
            if response.status_code == 200:
                data = response.json()
                if data.get("data") and len(data["data"]) > 0:
                    return data["data"][0]
            return None
        except Exception as e:
            console.print(f"[red]Error finding product by name: {e}[/red]")
            return None
    
    def find_by_reference(self, reference_string: str) -> Optional[Dict]:
        """Find a product by referenceString in Strapi"""
        try:
            response = requests.get(
                f"{self.base_url}/api/medical-products",
                params={"filters[referenceString][$eq]": reference_string},
                headers=self.headers
            )
            if response.status_code == 200:
                data = response.json()
                if data.get("data") and len(data["data"]) > 0:
                    return data["data"][0]
            return None
        except Exception as e:
            console.print(f"[red]Error finding product by reference: {e}[/red]")
            return None
    
    def create_product(self, payload: Dict) -> Optional[int]:
        """Create a new product in Strapi"""
        try:
            response = requests.post(
                f"{self.base_url}/api/medical-products",
                json=payload,
                headers=self.headers
            )
            if response.status_code in (200, 201):
                result = response.json()
                return result["data"]["id"]
            else:
                console.print(f"[red]Error creating product: {response.status_code} {response.text}[/red]")
                return None
        except Exception as e:
            console.print(f"[red]Error creating product: {e}[/red]")
            return None
    
    def update_product(self, product_id: int, payload: Dict) -> bool:
        """Update an existing product in Strapi"""
        try:
            response = requests.put(
                f"{self.base_url}/api/medical-products/{product_id}",
                json=payload,
                headers=self.headers
            )
            return response.status_code in (200, 201)
        except Exception as e:
            console.print(f"[red]Error updating product: {e}[/red]")
            return False
    
    def delete_product(self, product_id: int) -> bool:
        """Delete a product from Strapi"""
        try:
            response = requests.delete(
                f"{self.base_url}/api/medical-products/{product_id}",
                headers=self.headers
            )
            return response.status_code in (200, 204)
        except Exception as e:
            console.print(f"[red]Error deleting product: {e}[/red]")
            return False

class DiffUtils:
    """Utility class for JSON diff operations"""
    
    @staticmethod
    def print_diff(existing: Dict, new: Dict, title: str = "JSON Diff") -> None:
        """Print a side-by-side JSON diff"""
        console.print(f"\n[bold blue]{title}[/bold blue]")
        
        # Create a table for side-by-side comparison
        table = Table(title=title)
        table.add_column("Field", style="cyan")
        table.add_column("Existing", style="red")
        table.add_column("New", style="green")
        
        # Compare all fields
        all_fields = set(existing.keys()) | set(new.keys())
        
        for field in sorted(all_fields):
            existing_val = existing.get(field, "N/A")
            new_val = new.get(field, "N/A")
            
            # Convert to string for display
            existing_str = str(existing_val)[:50] + "..." if len(str(existing_val)) > 50 else str(existing_val)
            new_str = str(new_val)[:50] + "..." if len(str(new_val)) > 50 else str(new_val)
            
            # Highlight differences
            if existing_val != new_val:
                table.add_row(field, existing_str, new_str)
        
        console.print(table)

class ProgressTracker:
    """Handles progress tracking and resumability"""
    
    def __init__(self, progress_file: str = "progress.json"):
        self.progress_file = progress_file
        self.progress: List[ProgressEntry] = self.load_progress()
    
    def load_progress(self) -> List[ProgressEntry]:
        """Load progress from file"""
        if os.path.exists(self.progress_file):
            try:
                with open(self.progress_file, 'r') as f:
                    data = json.load(f)
                    return [ProgressEntry(**entry) for entry in data]
            except Exception as e:
                console.print(f"[yellow]Warning: Could not load progress file: {e}[/yellow]")
        return []
    
    def save_progress(self) -> None:
        """Save progress to file"""
        try:
            with open(self.progress_file, 'w') as f:
                json.dump([asdict(entry) for entry in self.progress], f, indent=2)
        except Exception as e:
            console.print(f"[red]Error saving progress: {e}[/red]")
    
    def add_entry(self, entry: ProgressEntry) -> None:
        """Add a progress entry and save"""
        self.progress.append(entry)
        self.save_progress()
    
    def get_processed_slides(self) -> List[int]:
        """Get list of processed slide numbers"""
        return [entry.slide_number for entry in self.progress]
    
    def is_slide_processed(self, slide_number: int) -> bool:
        """Check if a slide has been processed"""
        return slide_number in self.get_processed_slides()

def quick_text_guess(slide) -> str:
    """Extract a quick product name guess from slide text"""
    text_parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                text_parts.append(text)
    
    # Join all text and take first line as product name
    full_text = "\n".join(text_parts)
    if full_text:
        # Extract first line that looks like a product name
        lines = full_text.split('\n')
        for line in lines:
            line = line.strip()
            if line and not line.startswith('Slide') and len(line) > 3:
                # Remove slide numbers and clean up
                if '.' in line:
                    # Extract part after the reference number
                    parts = line.split(' ', 1)
                    if len(parts) > 1:
                        return parts[1]
                return line
    
    return "Unknown Product"

def extract_reference_string(slide) -> str:
    """Extract reference string from slide text"""
    text_parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                text_parts.append(text)
    
    # Join all text and look for reference pattern
    full_text = "\n".join(text_parts)
    if full_text:
        lines = full_text.split('\n')
        for line in lines:
            line = line.strip()
            if line and not line.startswith('Slide') and len(line) > 3:
                # Look for reference pattern (e.g., "4.1.1", "4.01.1")
                words = line.split()
                for word in words:
                    if '.' in word and word.replace('.', '').replace('0', '').replace('1', '').replace('2', '').replace('3', '').replace('4', '').replace('5', '').replace('6', '').replace('7', '').replace('8', '').replace('9', '') == '':
                        # This looks like a reference string
                        return word
    
    return ""

def extract_product_info(slide) -> tuple[str, str]:
    """Extract both product name and reference string from slide"""
    product_name = quick_text_guess(slide)
    reference_string = extract_reference_string(slide)
    return product_name, reference_string

def save_uploaded_product(product_id: int, payload: Dict, product_name: str) -> None:
    """Save uploaded product JSON to products_uploaded directory"""
    import os
    from datetime import datetime
    
    # Create products_uploaded directory if it doesn't exist
    upload_dir = "products_uploaded"
    os.makedirs(upload_dir, exist_ok=True)
    
    # Create filename with timestamp and product info
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = "".join(c for c in product_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
    safe_name = safe_name.replace(' ', '_')
    filename = f"{upload_dir}/{timestamp}_{product_id}_{safe_name}.json"
    
    # Add metadata to the payload
    upload_record = {
        "upload_metadata": {
            "uploaded_at": datetime.now().isoformat(),
            "product_id": product_id,
            "product_name": product_name,
            "upload_method": "cli_slide_driver"
        },
        "product_data": payload
    }
    
    # Save to file
    try:
        with open(filename, 'w', encoding='utf-8') as f:
            json.dump(upload_record, f, indent=2, ensure_ascii=False)
        console.print(f"[green]✓ Saved product data to: {filename}[/green]")
    except Exception as e:
        console.print(f"[red]Error saving product data: {e}[/red]")

def open_product_links(product_id: int) -> None:
    """Display both frontend and backend links for the product"""
    frontend_url = f"http://localhost:3000/products/{product_id}"
    backend_url = f"http://localhost:1337/admin/content-manager/collectionType/api::medical-product.medical-product/{product_id}?plugins[i18n][locale]=en"
    
    console.print(f"\n[bold green]Product uploaded successfully! (ID: {product_id})[/bold green]")
    console.print(f"[blue]Product links:[/blue]")
    console.print(f"[green]Frontend: {frontend_url}[/green]")
    console.print(f"[green]Backend:  {backend_url}[/green]")

def create_slide_banner(slide_info: SlideInfo) -> Panel:
    """Create a rich banner for slide information"""
    title = f"Slide {slide_info.slide_number}: {slide_info.product_name}"
    
    status_color = "green" if slide_info.existing_in_strapi else "yellow"
    status_text = "EXISTS" if slide_info.existing_in_strapi else "NEW"
    
    content = f"""
[bold]{slide_info.product_name}[/bold]
Reference: [cyan]{slide_info.reference_string}[/cyan]
Status: [{status_color}]{status_text}[/{status_color}]
"""
    
    if slide_info.strapi_id:
        content += f"Strapi ID: {slide_info.strapi_id}\n"
        content += f"URL: http://localhost:1337/admin/content-manager/collectionType/api::medical-product.medical-product/{slide_info.strapi_id}\n"
    
    return Panel(content, title=title, border_style="blue")

def process_slide_interactive(
    prs: Presentation,
    slide_index: int,
    strapi_utils: StrapiUtils,
    progress_tracker: ProgressTracker
) -> None:
    """Process a single slide interactively"""
    
    slide = prs.slides[slide_index]
    slide_number = slide_index + 1
    
    # Extract both product name and reference string
    product_name, reference_string = extract_product_info(slide)
    
    # Check if product exists in Strapi by reference string
    existing_product = strapi_utils.find_by_reference(reference_string) if reference_string else None
    existing_in_strapi = existing_product is not None
    strapi_id = existing_product["id"] if existing_product else None
    
    # Create slide info
    slide_info = SlideInfo(
        slide_number=slide_number,
        product_name=product_name,
        reference_string=reference_string,
        existing_in_strapi=existing_in_strapi,
        strapi_id=strapi_id
    )
    
    # Display banner
    console.print(create_slide_banner(slide_info))
    
    # Show action options
    console.print("\n[bold]Actions:[/bold]")
    console.print("A - Extract and show JSON")
    console.print("U - Upload to Strapi")
    console.print("D - Delete from Strapi")
    console.print("R - Delete and recreate")
    console.print("S - Skip this slide")
    console.print("Q - Quit")
    
    # Get user choice
    while True:
        try:
            choice = Prompt.ask("\nChoose action", choices=["A", "U", "D", "R", "S", "Q"])
            
            if choice == "Q":
                console.print("[yellow]Quitting...[/yellow]")
                sys.exit(0)
            
            elif choice == "S":
                progress_tracker.add_entry(ProgressEntry(
                    slide_number=slide_number,
                    action="skip",
                    timestamp=datetime.now().isoformat(),
                    product_name=product_name
                ))
                console.print("[yellow]Skipped[/yellow]")
                break
            
            elif choice == "A":
                # Extract and show JSON
                try:
                    payload = process_slide_two_steps(prs, slide_index)
                    console.print("[green]Extraction completed[/green]")
                    
                    # Show diff if product exists
                    if existing_in_strapi and existing_product:
                        existing_data = existing_product["attributes"]
                        new_data = payload["data"]
                        DiffUtils.print_diff(existing_data, new_data, f"Diff for {product_name}")
                    
                    # Ask user if they want to upload immediately
                    while True:
                        yn = input("Upload this JSON to Strapi now? [y/n]: ").strip().lower()
                        if yn in ("y", "n"):
                            break
                        print("Please enter 'y' or 'n'.")

                    if yn == "y":
                        # reuse the same helper that option 'U' uses
                        upload_success = strapi_utils.create_product(payload)
                        if upload_success:
                            open_product_links(upload_success)
                            save_uploaded_product(upload_success, payload, product_name)
                            progress_tracker.add_entry(ProgressEntry(
                                slide_number=slide_number,
                                action="upload",
                                timestamp=datetime.now().isoformat(),
                                product_name=product_name,
                                strapi_id=upload_success
                            ))
                        else:
                            console.print("[red]Upload failed[/red]")
                            progress_tracker.add_entry(ProgressEntry(
                                slide_number=slide_number,
                                action="upload",
                                timestamp=datetime.now().isoformat(),
                                product_name=product_name,
                                error="Upload failed"
                            ))
                    else:
                        print("🛈 Not uploaded; returning to menu.")
                        progress_tracker.add_entry(ProgressEntry(
                            slide_number=slide_number,
                            action="extract",
                            timestamp=datetime.now().isoformat(),
                            product_name=product_name
                        ))
                    
                except Exception as e:
                    console.print(f"[red]Extraction error: {e}[/red]")
                    progress_tracker.add_entry(ProgressEntry(
                        slide_number=slide_number,
                        action="extract",
                        timestamp=datetime.now().isoformat(),
                        product_name=product_name,
                        error=str(e)
                    ))
                break
            
            elif choice == "U":
                # Upload to Strapi
                try:
                    payload = process_slide_two_steps(prs, slide_index)
                    product_id = strapi_utils.create_product(payload)
                    
                    if product_id:
                        open_product_links(product_id)
                        save_uploaded_product(product_id, payload, product_name)
                        progress_tracker.add_entry(ProgressEntry(
                            slide_number=slide_number,
                            action="upload",
                            timestamp=datetime.now().isoformat(),
                            product_name=product_name,
                            strapi_id=product_id
                        ))
                    else:
                        console.print("[red]Upload failed[/red]")
                        progress_tracker.add_entry(ProgressEntry(
                            slide_number=slide_number,
                            action="upload",
                            timestamp=datetime.now().isoformat(),
                            product_name=product_name,
                            error="Upload failed"
                        ))
                except Exception as e:
                    console.print(f"[red]Upload error: {e}[/red]")
                    progress_tracker.add_entry(ProgressEntry(
                        slide_number=slide_number,
                        action="upload",
                        timestamp=datetime.now().isoformat(),
                        product_name=product_name,
                        error=str(e)
                    ))
                break
            
            elif choice == "D":
                # Delete from Strapi
                if existing_in_strapi and strapi_id:
                    if Confirm.ask(f"Delete product '{product_name}' (ID: {strapi_id})?"):
                        if strapi_utils.delete_product(strapi_id):
                            console.print("[green]Deleted successfully[/green]")
                            progress_tracker.add_entry(ProgressEntry(
                                slide_number=slide_number,
                                action="delete",
                                timestamp=datetime.now().isoformat(),
                                product_name=product_name,
                                strapi_id=strapi_id
                            ))
                        else:
                            console.print("[red]Delete failed[/red]")
                else:
                    console.print("[yellow]Product not found in Strapi[/yellow]")
                break
            
            elif choice == "R":
                # Delete and recreate
                if existing_in_strapi and strapi_id:
                    if Confirm.ask(f"Delete and recreate product '{product_name}' (ID: {strapi_id})?"):
                        # Delete first
                        if strapi_utils.delete_product(strapi_id):
                            console.print("[green]Deleted successfully[/green]")
                            
                            # Then create new
                            try:
                                payload = process_slide_two_steps(prs, slide_index)
                                new_product_id = strapi_utils.create_product(payload)
                                
                                if new_product_id:
                                    open_product_links(new_product_id)
                                    save_uploaded_product(new_product_id, payload, product_name)
                                    progress_tracker.add_entry(ProgressEntry(
                                        slide_number=slide_number,
                                        action="delete_recreate",
                                        timestamp=datetime.now().isoformat(),
                                        product_name=product_name,
                                        strapi_id=new_product_id
                                    ))
                                else:
                                    console.print("[red]Recreation failed[/red]")
                            except Exception as e:
                                console.print(f"[red]Recreation error: {e}[/red]")
                        else:
                            console.print("[red]Delete failed[/red]")
                else:
                    console.print("[yellow]Product not found in Strapi[/yellow]")
                break
            
        except KeyboardInterrupt:
            console.print("\n[yellow]Interrupted by user[/yellow]")
            progress_tracker.save_progress()
            sys.exit(0)

def main():
    """Main CLI entry point"""
    parser = argparse.ArgumentParser(description="Interactive CLI for processing PowerPoint slides")
    parser.add_argument("--pptx", default="data/WEB MASTER Ver 9.pptx", help="Path to PowerPoint file")
    parser.add_argument("--start", type=int, default=1, help="Starting slide number (1-based)")
    parser.add_argument("--end", type=int, help="Ending slide number (1-based)")
    parser.add_argument("--progress", default="progress.json", help="Progress file path")
    
    args = parser.parse_args()
    
    # Validate PowerPoint file
    if not os.path.exists(args.pptx):
        console.print(f"[red]PowerPoint file not found: {args.pptx}[/red]")
        sys.exit(1)
    
    # Initialize components
    strapi_utils = StrapiUtils()
    progress_tracker = ProgressTracker(args.progress)
    
    # Load PowerPoint
    try:
        prs = Presentation(args.pptx)
        total_slides = len(prs.slides)
        console.print(f"[green]Loaded presentation with {total_slides} slides[/green]")
    except Exception as e:
        console.print(f"[red]Error loading PowerPoint: {e}[/red]")
        sys.exit(1)
    
    # Determine slide range
    start_slide = max(1, args.start)
    end_slide = min(total_slides, args.end) if args.end else total_slides
    
    console.print(f"[blue]Processing slides {start_slide} to {end_slide}[/blue]")
    
    # Check for resumability
    processed_slides = progress_tracker.get_processed_slides()
    if processed_slides:
        console.print(f"[yellow]Found progress file with {len(processed_slides)} processed slides[/yellow]")
        if Confirm.ask("Resume from where you left off?"):
            # Skip already processed slides
            start_slide = max(start_slide, max(processed_slides) + 1)
            console.print(f"[blue]Resuming from slide {start_slide}[/blue]")
    
    # Process slides
    try:
        for slide_index in range(start_slide - 1, end_slide):
            console.print(f"\n{'='*60}")
            process_slide_interactive(prs, slide_index, strapi_utils, progress_tracker)
        
        console.print(f"\n[green]Completed processing slides {start_slide} to {end_slide}[/green]")
        
    except KeyboardInterrupt:
        console.print("\n[yellow]Interrupted by user. Progress saved.[/yellow]")
        progress_tracker.save_progress()
        sys.exit(0)

if __name__ == "__main__":
    main() 